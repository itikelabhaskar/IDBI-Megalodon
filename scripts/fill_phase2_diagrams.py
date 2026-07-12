"""Phase 2: officer workflow (slide 5) + prototype-now / sandbox-later architecture (slide 7)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"

# Bank-safe palette (IDBI-ish green, not purple AI-slop)
GREEN = RGBColor(0x1B, 0x5E, 0x3B)
GREEN_SOFT = RGBColor(0xE8, 0xF2, 0xEC)
GREEN_MID = RGBColor(0x2F, 0x7A, 0x52)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)  # terracotta for human-in-loop steps
ACCENT_SOFT = RGBColor(0xFB, 0xF0, 0xE8)
INK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0x6B, 0x7C, 0x72)


def clear_body(slide, min_top_inches: float = 1.2) -> None:
    to_del = []
    for shape in slide.shapes:
        # keep header picture (top), footer picture, and title text box near 0.8"
        if shape.shape_type == 13:  # picture
            top_in = shape.top / 914400
            if 0.5 < top_in < 5.4:  # shouldn't happen; keep header/footer
                continue
            continue
        if shape.has_text_frame and shape.shape_type == 17:
            top_in = shape.top / 914400
            if top_in > min_top_inches:
                to_del.append(shape)
        elif shape.shape_type not in (13, 17):
            # auto shapes / connectors we may have added
            top_in = shape.top / 914400
            if top_in > min_top_inches:
                to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)


def style_box(shape, fill, font_color=INK, font_size=10, bold=True) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.font.bold = bold
        p.font.color.rgb = font_color
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            run.font.bold = bold
            run.font.color.rgb = font_color


def add_box(slide, left, top, width, height, text, fill=GREEN_SOFT, font_size=10):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    # tighter corner
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass
    shape.text_frame.text = text
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    style_box(shape, fill, font_size=font_size)
    # vertical center-ish via anchor
    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow_right(slide, x1, y, x2, color=LINE):
    """Simple rightward chevron as a thin right-arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y - 0.08), Inches(x2 - x1), Inches(0.16)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, x, y1, y2, color=LINE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(x - 0.08), Inches(y1), Inches(0.16), Inches(y2 - y1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_caption(slide, text, left=0.3, top=5.05, width=9.4, font_size=10):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = "Calibri"
    p.font.italic = True
    p.font.color.rgb = MUTED
    return box


def build_slide5(slide) -> None:
    """Officer workflow with maker–checker — not a backend-only pipeline."""
    clear_body(slide)

    # Shrink title box slightly if needed — leave template title text
    for shape in slide.shapes:
        if shape.has_text_frame and "Process flow" in shape.text_frame.text:
            shape.text_frame.paragraphs[0].font.size = Pt(16)
            shape.text_frame.paragraphs[0].font.bold = True

    # Row 1: engine path
    row1 = [
        (0.25, "Consent"),
        (1.65, "Source\nfetch"),
        (3.05, "HealthScore\n+ BRE"),
        (4.55, "Go /\nConditional /\nNo-Go"),
    ]
    box_w, box_h = 1.25, 0.85
    y1 = 1.55
    for i, (x, label) in enumerate(row1):
        add_box(slide, x, y1, box_w, box_h, label, GREEN_SOFT, font_size=10)
        if i < len(row1) - 1:
            add_arrow_right(slide, x + box_w + 0.02, y1 + box_h / 2, row1[i + 1][0] - 0.02)

    # Down from Go/No-Go to Maker
    add_arrow_down(slide, 5.15, y1 + box_h + 0.05, 2.75)

    # Row 2: human-in-loop (accent)
    row2 = [
        (4.55, "Maker\nreview", ACCENT_SOFT),
        (6.15, "Checker\nsanction /\nreturn", ACCENT_SOFT),
        (7.85, "CAM +\naudit", GREEN_SOFT),
    ]
    y2 = 2.85
    for i, (x, label, fill) in enumerate(row2):
        add_box(slide, x, y2, box_w + 0.15, box_h, label, fill, font_size=10)
        if i < len(row2) - 1:
            add_arrow_right(
                slide,
                x + box_w + 0.17,
                y2 + box_h / 2,
                row2[i + 1][0] - 0.02,
                ACCENT if i == 0 else LINE,
            )

    # Left-side note: engine vs human
    note = slide.shapes.add_textbox(Inches(0.25), Inches(2.85), Inches(4.0), Inches(1.1))
    tf = note.text_frame
    tf.word_wrap = True
    lines = [
        ("Engine path", True, GREEN),
        ("Consented data → HealthScore + BRE → Go / Conditional / No-Go", False, MUTED),
        ("", False, MUTED),
        ("Human path (AMA)", True, ACCENT),
        ("Credit Officer = maker · Risk Admin = checker", False, MUTED),
        ("Recommendation ≠ final sanction", False, MUTED),
    ]
    for i, (t, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(10 if bold else 9)
        p.font.name = "Calibri"
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(1)

    add_caption(
        slide,
        "The engine recommends. The IDBI underwriter decides. Maker–checker keeps human intervention in the loop.",
        top=4.95,
    )


def build_slide7(slide) -> None:
    """Prototype now vs sandbox later — same frozen MsmeCase contract.

    Layout rule: left and right panels must NOT be overlapped by the bridge.
    Gap between panels is ~1.1\" (left ends 4.45, right starts 5.55); the
    MsmeCase callout lives only inside that gap as a narrow vertical bridge.
    """
    clear_body(slide)

    for shape in slide.shapes:
        if shape.has_text_frame and "Architecture" in shape.text_frame.text:
            shape.text_frame.paragraphs[0].font.size = Pt(16)
            shape.text_frame.paragraphs[0].font.bold = True

    # Two panels with a clean center gap (no overlap with bridge)
    # Left:  0.25 → 4.45  (width 4.20)
    # Gap:   4.45 → 5.55  (width 1.10)  ← bridge lives here only
    # Right: 5.55 → 9.75  (width 4.20)
    panel_h = 3.15
    panel_top = 1.45
    left_w = 4.20
    right_w = 4.20
    left_x = 0.25
    right_x = 5.55

    left = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_x), Inches(panel_top), Inches(left_w), Inches(panel_h)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = GREEN_SOFT
    left.line.color.rgb = GREEN
    left.line.width = Pt(1.5)
    try:
        left.adjustments[0] = 0.08
    except Exception:
        pass

    right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x),
        Inches(panel_top),
        Inches(right_w),
        Inches(panel_h),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF5)
    right.line.color.rgb = GREEN
    right.line.width = Pt(1.5)
    try:
        right.adjustments[0] = 0.08
    except Exception:
        pass

    def fill_panel(shape, title, bullets, title_color=GREEN):
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.name = "Calibri"
        p0.font.color.rgb = title_color
        p0.space_after = Pt(8)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            p.font.color.rgb = INK
            p.space_after = Pt(4)
            p.level = 0

    fill_panel(
        left,
        "NOW — prototype",
        [
            "• Synthetic MSME population (deterministic)",
            "• Connector stubs: AA / GSTN / UPI / EPFO /",
            "  DISCOM power / fuel / bureau-lite",
            "• Features → HealthScore + BRE (+ ML viability)",
            "• Officer UI + maker–checker + CAM",
            "• ULI / OCEN / AA-style contract stubs",
            "  (Zod-validated; not live REST rails)",
        ],
    )

    fill_panel(
        right,
        "LATER — after shortlist",
        [
            "• Same frozen MsmeCase contract",
            "• Swap stubs for IDBI / partner sandbox APIs",
            "• Recalibrate on anonymised bank data",
            "• Approved cloud / IDBI sandbox path",
            "  for pilot hosting",
            "• Officer UAT with bank maker–checker",
            "• Harden consent artifacts + audit export",
        ],
        title_color=ACCENT,
    )

    # Narrow vertical bridge ONLY in the center gap (x 4.52–5.48)
    bridge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.52), Inches(1.95), Inches(0.96), Inches(2.15)
    )
    bridge.fill.solid()
    bridge.fill.fore_color.rgb = WHITE
    bridge.line.color.rgb = ACCENT
    bridge.line.width = Pt(1.5)
    try:
        bridge.adjustments[0] = 0.15
    except Exception:
        pass
    bridge.text_frame.clear()
    bridge.text_frame.word_wrap = True
    bridge.text_frame.anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(["Same", "MsmeCase", "field names", "frozen"]):
        p = bridge.text_frame.paragraphs[0] if i == 0 else bridge.text_frame.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = ACCENT
        p.space_after = Pt(2)

    add_caption(
        slide,
        "Not live ULI/OCEN/AA integration — sandbox-ready integration surface that swaps without rewriting the officer UI.",
        top=4.95,
    )


def main() -> None:
    prs = Presentation(PATH)
    # Slide 5 already approved — only rebuild Slide 7
    build_slide7(prs.slides[6])
    prs.save(PATH)
    print("SAVED", PATH)

    # Geometry proof: bridge must sit strictly between left end and right start
    slide = prs.slides[6]
    panels = []
    bridge = None
    for shape in slide.shapes:
        if shape.shape_type != 1:  # AUTO_SHAPE
            continue
        left = shape.left / 914400
        top = shape.top / 914400
        width = shape.width / 914400
        height = shape.height / 914400
        text = shape.text_frame.text if shape.has_text_frame else ""
        if "NOW" in text:
            panels.append(("LEFT", left, left + width, top, top + height))
        elif "LATER" in text:
            panels.append(("RIGHT", left, left + width, top, top + height))
        elif "MsmeCase" in text:
            bridge = ("BRIDGE", left, left + width, top, top + height, text.replace("\n", " | "))

    for name, x0, x1, y0, y1 in panels:
        print(f"{name}: x {x0:.2f}–{x1:.2f}, y {y0:.2f}–{y1:.2f}")
    if bridge:
        name, x0, x1, y0, y1, t = bridge
        print(f"{name}: x {x0:.2f}–{x1:.2f}, y {y0:.2f}–{y1:.2f}  [{t}]")
        left = next(p for p in panels if p[0] == "LEFT")
        right = next(p for p in panels if p[0] == "RIGHT")
        overlaps_left = x0 < left[2] - 0.01 and x1 > left[1] + 0.01
        overlaps_right = x0 < right[2] - 0.01 and x1 > right[1] + 0.01
        print(f"overlaps_left={overlaps_left} overlaps_right={overlaps_right}")
        assert not overlaps_left and not overlaps_right, "Bridge still overlaps a column"


if __name__ == "__main__":
    main()
