"""Phase 3: selective screenshots on slides 6 (officer workflow) and 10 (model/governance).

Friend guidance: selective, not a tiny collage — 2–3 readable shots max.
Captions tied to judging / AMA criteria.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"
SHOTS = r"D:\IDBI-Megalodon\archive\ppt-screenshots"

GREEN = RGBColor(0x1B, 0x5E, 0x3B)
MUTED = RGBColor(0x55, 0x55, 0x55)
SOFT = RGBColor(0xE8, 0xF2, 0xEC)


def clear_body(slide, min_top_inches: float = 1.35) -> None:
    """Remove prior body content below the template title; keep header/footer pictures."""
    to_del = []
    for shape in slide.shapes:
        top_in = shape.top / 914400
        if top_in <= min_top_inches:
            continue
        if shape.shape_type == 13 and top_in >= 5.4:
            continue
        to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)


def add_caption(slide, left, top, width, height, text: str, size: int = 9) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(size)
    p.font.name = "Calibri"
    p.font.color.rgb = MUTED


def add_label(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    try:
        shape.adjustments[0] = 0.2
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.color.rgb = GREEN


def add_shot(slide, path: str, left, top, width, max_height: float = 3.15) -> None:
    """Fit by width, then cap height so captions stay above the footer."""
    from PIL import Image

    with Image.open(path) as im:
        w_px, h_px = im.size
    height = width * (h_px / w_px)
    if height > max_height:
        # re-add with height cap (width shrinks to preserve aspect)
        height = max_height
        width = height * (w_px / h_px)
        # keep left edge; slight right-align unused space is fine
    slide.shapes.add_picture(
        path, Inches(left), Inches(top), width=Inches(width), height=Inches(height)
    )


def fill_slide_6(prs: Presentation) -> None:
    """Officer workflow: Health Card (Go) + Maker–Checker. Queue as third readable strip."""
    slide = prs.slides[5]
    clear_body(slide)

    # Layout: two large proof shots (readable) + one narrower queue strip
    # Content band ~1.40 → 5.40 (footer at 5.52)
    label_top = 1.40
    label_h = 0.24
    img_top = 1.68

    # Left: Health Card (primary workflow proof)
    add_label(slide, 0.25, label_top, 4.55, label_h, "1 · Health Card — Go / No-Go")
    add_shot(slide, f"{SHOTS}\\02-healthcard-go-crop.png", 0.25, img_top, 4.55, max_height=3.15)
    add_caption(
        slide,
        0.25,
        4.95,
        4.55,
        0.48,
        "Underwriter view: Go recommendation, HealthScore, limit. Officer owns the decision.",
        size=9,
    )

    # Right: Maker–Checker
    add_label(slide, 5.00, label_top, 4.75, label_h, "2 · Maker → Checker (four-eyes)")
    add_shot(slide, f"{SHOTS}\\03-decision-maker-checker-crop.png", 5.00, img_top, 4.75, max_height=3.15)
    add_caption(
        slide,
        5.00,
        4.95,
        4.75,
        0.48,
        "Engine recommends; maker records Approve/Refer/Reject; checker sanctions. Not auto-sanction.",
        size=9,
    )


def fill_slide_10(prs: Presentation) -> None:
    """Model / governance: model card + champion–challenger (2 large readable shots)."""
    slide = prs.slides[9]
    clear_body(slide)

    label_top = 1.40
    label_h = 0.24
    img_top = 1.68

    add_label(slide, 0.25, label_top, 4.55, label_h, "1 · Model card (viability scorecard)")
    add_shot(slide, f"{SHOTS}\\05-governance-model-crop.png", 0.25, img_top, 4.55, max_height=3.15)
    add_caption(
        slide,
        0.25,
        4.95,
        4.55,
        0.48,
        "EPFO + Power in inputs · held-out AUC 0.745 · advisory tier — officer accepts/overrides · sandbox recalibrate.",
        size=9,
    )

    add_label(slide, 5.00, label_top, 4.75, label_h, "2 · Champion–challenger & calibration")
    add_shot(slide, f"{SHOTS}\\06-champion-challenger-crop.png", 5.00, img_top, 4.75, max_height=3.15)
    add_caption(
        slide,
        5.00,
        4.95,
        4.75,
        0.48,
        "Explainable logistic kept as control; GBM challenger within noise — governance over black-box chase.",
        size=9,
    )


def main() -> None:
    prs = Presentation(PATH)
    fill_slide_6(prs)
    fill_slide_10(prs)
    prs.save(PATH)
    print("Phase 3 done: slides 6 + 10 — 2 readable screenshots each.")


if __name__ == "__main__":
    main()
