"""Phase 5: slides 14–15 (why shortlist + thank you) + AMA checklist verification."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"

GREEN = RGBColor(0x1B, 0x5E, 0x3B)
GREEN_SOFT = RGBColor(0xE8, 0xF2, 0xEC)
INK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC5, 0xD0, 0xC9)
LINK = RGBColor(0x05, 0x63, 0xC1)


def clear_body(slide, min_top_inches: float = 0.7, keep_footer: bool = True) -> None:
    to_del = []
    for shape in slide.shapes:
        top_in = shape.top / 914400
        # keep header strip
        if shape.shape_type == 13 and top_in < 0.7:
            continue
        if keep_footer and shape.shape_type == 13 and top_in >= 5.4:
            continue
        # keep full-bleed thank-you background on slide 15
        if shape.shape_type == 13 and shape.height / 914400 > 5.0:
            continue
        if top_in < min_top_inches and shape.shape_type == 13:
            continue
        to_del.append(shape)
    for shape in to_del:
        # don't delete placeholders on thank-you — we fill them
        if shape.has_text_frame and shape.is_placeholder:
            continue
        el = shape._element
        el.getparent().remove(el)


def set_run(p, text, size=12, bold=False, color=INK, name="Calibri"):
    # clear existing runs
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = color
    return run


def fill_slide_14(prs: Presentation) -> None:
    """Why shortlist — map to AMA checklist. Fewer claims, stronger proof."""
    slide = prs.slides[13]
    # Remove any prior body (keep header/footer pics)
    to_del = []
    for shape in slide.shapes:
        top_in = shape.top / 914400
        if shape.shape_type == 13 and (top_in < 0.7 or top_in >= 5.4):
            continue
        to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)

    # Title
    title = slide.shapes.add_textbox(Inches(0.30), Inches(0.85), Inches(9.4), Inches(0.40))
    tf = title.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0], "Why shortlist HealthLens", size=18, bold=True, color=INK)

    sub = slide.shapes.add_textbox(Inches(0.30), Inches(1.25), Inches(9.4), Inches(0.35))
    set_run(
        sub.text_frame.paragraphs[0],
        "Maps directly to Track 03 AMA asks — decision support, not auto-sanction.",
        size=11,
        bold=False,
        color=MUTED,
    )

    reasons = [
        (
            "1",
            "Power + fuel authenticity",
            "Electricity and fuel sit in operations with sector authenticity bands — AMA signals, not chart decoration.",
        ),
        (
            "2",
            "Go / No-Go · underwriter decides",
            "BRE recommends Approve / Refer / Reject / Incomplete; Low confidence forces maker–checker. AI gives reason only.",
        ),
        (
            "3",
            "Thin-file NTC/NTB not auto-rejected",
            "Credit-invisible MSMEs stay in the funnel; Credit-Invisible Lift + scheme readiness on the workbench.",
        ),
        (
            "4",
            "Pilot-ready path, honest today",
            "Live prototype + rails honesty panel + SETUP_GAPS → IDBI sandbox APIs on an approved AWS path after shortlist.",
        ),
    ]

    for i, (num, headline, body) in enumerate(reasons):
        top = 1.70 + i * 0.88
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.30), Inches(top + 0.08), Inches(0.36), Inches(0.36)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        btf = badge.text_frame
        btf.clear()
        bp = btf.paragraphs[0]
        set_run(bp, num, size=12, bold=True, color=WHITE)
        bp.alignment = PP_ALIGN.CENTER

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.80), Inches(top), Inches(8.90), Inches(0.78)
        )
        try:
            card.adjustments[0] = 0.08
        except Exception:
            pass
        card.fill.solid()
        card.fill.fore_color.rgb = GREEN_SOFT
        card.line.color.rgb = LINE
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.clear()
        p0 = ctf.paragraphs[0]
        set_run(p0, headline, size=12, bold=True, color=GREEN)
        p0.space_after = Pt(2)
        p1 = ctf.add_paragraph()
        set_run(p1, body, size=10, bold=False, color=INK)


def fill_slide_15(prs: Presentation) -> None:
    """Thank you — team + product URL (clickable)."""
    slide = prs.slides[14]

    # Fill placeholders if present; else add text boxes
    placeholders = [sh for sh in slide.shapes if sh.has_text_frame and sh.is_placeholder]
    placeholders.sort(key=lambda s: s.top)

    thank_lines = [
        ("Thank you", 32, True, WHITE if _has_dark_bg(slide) else INK),
        ("Team Megalodon  ·  Leader: Bhaskar Itikela", 16, False, None),
        ("IDBI MSME HealthLens — Track 03", 14, False, None),
    ]

    # Prefer writing into existing placeholders
    if len(placeholders) >= 1:
        tf = placeholders[0].text_frame
        tf.clear()
        tf.word_wrap = True
        # Detect text color: slide 15 has full-bleed image — use white
        color = WHITE
        p0 = tf.paragraphs[0]
        set_run(p0, "Thank you", size=36, bold=True, color=color)
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(10)
        p1 = tf.add_paragraph()
        set_run(p1, "Team Megalodon  ·  Leader: Bhaskar Itikela", size=16, bold=False, color=color)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(6)
        p2 = tf.add_paragraph()
        set_run(p2, "IDBI MSME HealthLens — Track 03", size=14, bold=False, color=color)
        p2.alignment = PP_ALIGN.CENTER

    if len(placeholders) >= 2:
        tf = placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        color = WHITE
        p0 = tf.paragraphs[0]
        run = set_run(p0, "https://idbi-megalodon.vercel.app", size=14, bold=False, color=LINK)
        run.hyperlink.address = "https://idbi-megalodon.vercel.app"
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(4)
        p1 = tf.add_paragraph()
        set_run(
            p1,
            "Live prototype on synthetic data  ·  GitHub: itikelabhaskar/IDBI-Megalodon",
            size=11,
            bold=False,
            color=color,
        )
        p1.alignment = PP_ALIGN.CENTER
    else:
        # Fallback text boxes
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0], "Thank you", size=36, bold=True, color=INK)


def _has_dark_bg(slide) -> bool:
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.height / 914400 > 5.0:
            return True
    return False


def verify_ama_checklist(prs: Presentation) -> list[str]:
    """Return missing AMA signals (empty = all present somewhere in deck)."""
    corpus = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                corpus.append(shape.text_frame.text.lower())
    text = "\n".join(corpus)

    checks = [
        ("A alternate/NTC", any(k in text for k in ("ntc", "ntb", "alternate", "credit-invisible"))),
        ("B power/electricity", any(k in text for k in ("power", "electricity"))),
        ("C EPFO", "epfo" in text),
        ("D Go/No-Go", "go" in text and ("no-go" in text or "nogo" in text or "no go" in text)),
        ("E fuel + thin-file", "fuel" in text and ("thin" in text or "ntc" in text)),
        ("F underwriter/maker", any(k in text for k in ("underwriter", "maker", "officer owns"))),
        ("G sandbox", "sandbox" in text),
        ("H AWS", "aws" in text),
        ("I stubs ULI/OCEN/AA", "stub" in text and ("uli" in text or "ocen" in text or "aa" in text)),
    ]
    missing = [name for name, ok in checks if not ok]
    return missing


def main() -> None:
    prs = Presentation(PATH)
    fill_slide_14(prs)
    fill_slide_15(prs)
    prs.save(PATH)

    # re-open and verify
    prs2 = Presentation(PATH)
    missing = verify_ama_checklist(prs2)
    print("Phase 5 done: slides 14–15.")
    if missing:
        print("AMA gaps:", missing)
    else:
        print("AMA checklist: all rows present in deck copy.")

    # hyperlink count on slides 13 + 15
    from pptx.oxml.ns import qn

    for idx in (12, 14):
        slide = prs2.slides[idx]
        xml = slide._element.xml
        count = xml.count("hlinkClick")
        print(f"Slide {idx + 1} hlinkClick count: {count}")


if __name__ == "__main__":
    main()
