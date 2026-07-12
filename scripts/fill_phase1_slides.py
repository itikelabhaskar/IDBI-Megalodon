"""Phase 1: fill slides 1–4 (indices 0–3) of the IDBI Innovate submission deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"


def clear_body_textboxes(slide, min_top_inches: float = 1.2) -> None:
    """Remove previously added body text boxes below the template title area."""
    to_del = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            top_in = shape.top / 914400
            if top_in > min_top_inches:
                to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)


def add_body(
    slide,
    text: str,
    left: float = 0.3,
    top: float = 1.5,
    width: float = 9.4,
    height: float = 3.5,
    font_size: int = 13,
    hook_lines: int = 0,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        is_hook = i < hook_lines
        p.font.size = Pt(15 if is_hook else font_size)
        p.font.name = "Calibri"
        p.font.bold = is_hook
        p.font.color.rgb = RGBColor(0x14, 0x14, 0x14) if is_hook else RGBColor(0x28, 0x28, 0x28)
        p.space_after = Pt(6 if is_hook else 3)


def main() -> None:
    prs = Presentation(PATH)

    # Drop MCP-added body boxes on slides 1–3 so we don't double up
    for si in (1, 2, 3):
        clear_body_textboxes(prs.slides[si])

    # --- Slide 0: Team Details ---
    s0 = prs.slides[0]
    for shape in s0.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        content = [
            ("Team Details", 14, True),
            ("", 11, False),
            ("Team name: Megalodon", 13, False),
            ("Team leader: Bhaskar Itikela", 13, False),
            (
                "Problem Statement: Track 03 — Financial Inclusion / Digital Lending / Credit Decisioning",
                12,
                False,
            ),
            ("", 11, False),
            (
                "Product: IDBI MSME HealthLens — officer-ready Financial Health Card from consented alternate data.",
                12,
                False,
            ),
            ("Live prototype + GitHub ready for sandbox refinement.", 12, False),
        ]
        for i, (line, size, bold) in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.name = "Calibri"
            p.font.bold = bold
            p.font.color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
            p.space_after = Pt(2)

    # --- Slide 1: Brief ---
    add_body(
        prs.slides[1],
        "Balance sheets are thin. Business behaviour is not.\n"
        "HealthLens turns GST, bank, UPI, EPFO, power and fuel signals into an officer-ready MSME Health Card.\n"
        "\n"
        "• Credit-officer workbench: HealthScore, Go / Conditional / No-Go, limit + IDBI product route, "
        "with reasons an underwriter can defend.\n"
        "• Not auto-sanction — IDBI remains the decision owner; AI gives reason and logic only.\n"
        "• Thin-file NTC/NTB stay in the funnel: when GST or bureau is weak, power and fuel help build "
        "the turnover picture.\n"
        "• Proof now: deployed on synthetic, leakage-aware data; connectors swap to IDBI sandbox APIs "
        "after shortlist.",
        top=1.45,
        height=3.7,
        font_size=12,
        hook_lines=2,
    )

    # --- Slide 2: Opportunities ---
    add_body(
        prs.slides[2],
        "How different?\n"
        "• Full credit workflow — not a score widget: consent → HealthScore + BRE → maker–checker → CAM → governance.\n"
        "• Power (manufacturing) and fuel (traders/logistics) sit inside the operations sub-score, not only as charts.\n"
        "• Reason codes explain the decision; ML contributions explain only the secondary viability proxy.\n"
        "\n"
        "How it solves the problem?\n"
        "• Credit-invisible MSMEs are not auto-rejected when traditional papers are thin.\n"
        "• GST vs bank vs power/fuel triangulation catches inflated turnover.\n"
        "• Routes to real IDBI products (i-MSME Express, GeM Sahay, CGTMSE) with eligibility gates.\n"
        "\n"
        "USP: Pilot-ready decision-support workbench — alternate data → explainable Health Card → "
        "officer Go/No-Go — ULI/OCEN/AA-style stubs ready for sandbox swap.\n"
        "\n"
        "Proof: Frozen MsmeCase contract · deterministic generator · leakage guard · audit hash chain · 93 tests",
        top=2.05,
        height=3.3,
        font_size=11,
        hook_lines=0,
    )

    # --- Slide 3: Features ---
    add_body(
        prs.slides[3],
        "• MSME Financial Health Card — 7 sub-scores → HealthScore + A–D band "
        "(incl. EPFO, power / fuel operations)\n"
        "• Go / Conditional / No-Go with recommended limit — speaks officer language\n"
        "• Business-need detection + IDBI product routing (GeM Sahay gates included)\n"
        "• Fraud & triangulation — GST–bank, power, fuel; Benford advisory only\n"
        "• Explainability + path-to-credit for Refer / Reject — underwriter keeps the call\n"
        "• Maker–checker workflow + printable CAM preview\n"
        "• Live applicant scoring + governance (model card, champion–challenger, audit trail)",
        top=1.45,
        height=3.7,
        font_size=13,
        hook_lines=0,
    )

    prs.save(PATH)
    print("SAVED", PATH)

    # Verify
    prs2 = Presentation(PATH)
    for si in range(4):
        print(f"=== SLIDE {si + 1} ===")
        for shape in prs2.slides[si].shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print(shape.text_frame.text[:400])
                print("---")


if __name__ == "__main__":
    main()
