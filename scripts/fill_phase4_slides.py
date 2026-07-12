"""Phase 4: slides 8–13 (tech, effort, metrics, roadmap, links) + optional Slide 6 title polish.

Standard: fewer claims, stronger proof. No live ULI/OCEN/AA. No production-ready claims.
Demo video left blank.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"

GREEN = RGBColor(0x1B, 0x5E, 0x3B)
GREEN_SOFT = RGBColor(0xE8, 0xF2, 0xEC)
INK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC5, 0xD0, 0xC9)


def clear_body(slide, min_top_inches: float = 1.35) -> None:
    to_del = []
    for shape in slide.shapes:
        top_in = shape.top / 914400
        if top_in <= min_top_inches:
            continue
        if shape.shape_type == 13 and top_in >= 5.4:
            continue
        to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)


def set_run(p, text, size=12, bold=False, color=INK, name="Calibri"):
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = color
    return run


def add_textbox(slide, left, top, width, height, lines, font_size=12):
    """lines: list of (text, size, bold, color) or plain str."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, font_size, False, INK
        else:
            text, size, bold, color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_run(p, text, size=size, bold=bold, color=color)
        p.space_after = Pt(4 if bold else 2)
        p.level = 0
    return box


def add_panel(slide, left, top, width, height, title, bullets, title_size=11, body_size=11):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_SOFT
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    # Title strip as text inside
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    set_run(p0, title, size=title_size, bold=True, color=GREEN)
    p0.space_after = Pt(6)
    for b in bullets:
        p = tf.add_paragraph()
        set_run(p, b, size=body_size, bold=False, color=INK)
        p.space_after = Pt(3)
    return shape


def polish_slide6_title(prs: Presentation) -> None:
    """Optional: template title is weak vs content. Editable text box — safe to retitle."""
    slide = prs.slides[5]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "Wireframes" in shape.text_frame.text or "Mock diagrams" in shape.text_frame.text:
            tf = shape.text_frame
            p = tf.paragraphs[0]
            if p.runs:
                p.runs[0].text = "Officer workflow in the prototype"
            else:
                set_run(p, "Officer workflow in the prototype", size=16, bold=True)
            return


def fill_slide_8(prs: Presentation) -> None:
    """Technologies — honest stack + AWS/sandbox portability. No live-integration claims."""
    slide = prs.slides[7]
    clear_body(slide)

    add_textbox(
        slide,
        0.3,
        1.40,
        9.4,
        0.35,
        [
            (
                "Built now on a portable stack — same codebase can target an IDBI-approved AWS path.",
                12,
                True,
                INK,
            )
        ],
    )

    # Three columns
    cols = [
        (
            0.25,
            "App & UI",
            [
                "• TanStack Start + React 19",
                "• TypeScript · Tailwind v4",
                "• shadcn/ui · Recharts",
                "• Role-based officer workbench",
            ],
        ),
        (
            3.45,
            "Scoring & contracts",
            [
                "• Zod-validated MsmeCase",
                "• Rule scorecard + BRE",
                "• Transparent ML proxy (secondary)",
                "• ULI/OCEN/AA-style stubs (not live)",
            ],
        ),
        (
            6.65,
            "Deploy & ops",
            [
                "• Nitro server (portable)",
                "• Live on Vercel for demo",
                "• Presets: AWS Lambda / Node",
                "• Vitest — 93 tests passing",
            ],
        ),
    ]
    for left, title, bullets in cols:
        add_panel(slide, left, 1.85, 3.05, 2.55, title, bullets, title_size=12, body_size=11)

    add_textbox(
        slide,
        0.3,
        4.55,
        9.4,
        0.75,
        [
            (
                "Honesty line: connectors are stubs with Zod contracts — swap to IDBI sandbox APIs after shortlist. "
                "Not claiming live AA / ULI / OCEN / GSTN / DISCOM today.",
                10,
                False,
                MUTED,
            ),
        ],
    )


def fill_slide_9(prs: Presentation) -> None:
    """Effort bands — no invented ₹ TCO. Built now vs sandbox pilot."""
    slide = prs.slides[8]
    clear_body(slide)

    add_textbox(
        slide,
        0.3,
        1.40,
        9.4,
        0.40,
        [
            (
                "Effort framing (not a commercial quote) — what is built vs what a sandbox pilot needs.",
                12,
                True,
                INK,
            )
        ],
    )

    add_panel(
        slide,
        0.25,
        1.95,
        4.60,
        2.85,
        "BUILT NOW (prototype)",
        [
            "• Full officer workflow + Health Card",
            "• Power + fuel in operations score",
            "• Maker–checker + CAM + audit trail",
            "• Governance: model card, champion–challenger",
            "• Synthetic book (600) + live demo deploy",
            "• Frozen contract ready for API swap",
        ],
        title_size=12,
        body_size=11,
    )

    add_panel(
        slide,
        5.15,
        1.95,
        4.60,
        2.85,
        "SANDBOX PILOT (after shortlist)",
        [
            "• Wire IDBI / partner sandbox APIs",
            "• Recalibrate on bank-labelled outcomes",
            "• Policy thresholds with risk owners",
            "• UAT with credit officers",
            "• Host on approved cloud / AWS path",
            "• Security, consent, and ops hardening",
        ],
        title_size=12,
        body_size=11,
    )


def fill_slide_11(prs: Presentation) -> None:
    """Exact metrics + synthetic caveat. Strong proof, no production AUC claim."""
    slide = prs.slides[10]
    clear_body(slide)

    add_textbox(
        slide,
        0.3,
        1.40,
        9.4,
        0.35,
        [("Prototype book — synthetic, leakage-aware. Figures are indicative until sandbox backtest.", 11, True, MUTED)],
    )

    # Metric tiles
    metrics = [
        ("600", "MSMEs scored", "Deterministic synthetic population"),
        ("376 / 112 / 112", "Approve · Refer · Reject", "BRE + HealthScore recommendation"),
        ("143 / 143", "NTC/NTB kept in funnel", "No auto-reject of credit-invisible"),
        ("0.7449", "Held-out AUC (logistic)", "KS 0.3609 · Gini 0.4898"),
        ("0.7498", "GBM challenger AUC", "Explainable champion kept as control"),
        ("93", "Automated tests", "Vitest — engine + BRE + UI contracts"),
    ]

    # 2 rows × 3
    for i, (value, label, sub) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left = 0.25 + col * 3.20
        top = 1.85 + row * 1.45
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.05), Inches(1.30)
        )
        try:
            shape.adjustments[0] = 0.1
        except Exception:
            pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LINE
        shape.line.width = Pt(1.25)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        set_run(p0, value, size=18, bold=True, color=GREEN)
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(2)
        p1 = tf.add_paragraph()
        set_run(p1, label, size=11, bold=True, color=INK)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(2)
        p2 = tf.add_paragraph()
        set_run(p2, sub, size=9, bold=False, color=MUTED)
        p2.alignment = PP_ALIGN.CENTER

    add_textbox(
        slide,
        0.3,
        4.85,
        9.4,
        0.50,
        [
            (
                "Caveat: AUC/KS are on synthetic held-out labels — not production performance. "
                "ML is a secondary viability proxy; HealthScore + BRE drive the officer recommendation.",
                9,
                False,
                MUTED,
            )
        ],
    )


def fill_slide_12(prs: Presentation) -> None:
    """Future development — sandbox → calibrate → harden → UAT."""
    slide = prs.slides[11]
    clear_body(slide)

    add_textbox(
        slide,
        0.3,
        1.40,
        9.4,
        0.35,
        [("After shortlist — same frozen MsmeCase, real sandbox data, bank-owned thresholds.", 12, True, INK)],
    )

    steps = [
        ("1", "Sandbox APIs", "Swap stubs for IDBI / partner feeds (AA, GSTN, UPI, EPFO, power)"),
        ("2", "Recalibrate", "Refit scorecard + ML on bank-labelled outcomes; freeze policy bands"),
        ("3", "Harden", "Consent, audit, RBAC, observability on approved cloud / AWS path"),
        ("4", "UAT", "Credit-officer walkthroughs; maker–checker in bank process"),
    ]

    for i, (num, title, body) in enumerate(steps):
        left = 0.25 + i * 2.40
        # number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + 0.85), Inches(1.90), Inches(0.40), Inches(0.40)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        set_run(p, num, size=12, bold=True, color=WHITE)
        p.alignment = PP_ALIGN.CENTER

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.45), Inches(2.25), Inches(2.20)
        )
        try:
            card.adjustments[0] = 0.1
        except Exception:
            pass
        card.fill.solid()
        card.fill.fore_color.rgb = GREEN_SOFT
        card.line.color.rgb = LINE
        tf = card.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        set_run(p0, title, size=12, bold=True, color=GREEN)
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(6)
        p1 = tf.add_paragraph()
        set_run(p1, body, size=10, bold=False, color=INK)
        p1.alignment = PP_ALIGN.CENTER

    add_textbox(
        slide,
        0.3,
        4.85,
        9.4,
        0.50,
        [
            (
                "Out of scope until bank engagement: live production sanctioning, guaranteed AUC, "
                "or claiming ULI/OCEN/AA are already integrated.",
                9,
                False,
                MUTED,
            )
        ],
    )


def fill_slide_13(prs: Presentation) -> None:
    """Links — GitHub + product; demo video blank."""
    slide = prs.slides[12]
    clear_body(slide, min_top_inches=2.2)  # keep template prompt block

    # Soften / keep template heading; add filled link cards below
    add_textbox(
        slide,
        0.3,
        2.35,
        9.4,
        0.30,
        [("Submission links", 12, True, GREEN)],
    )

    links = [
        (
            "GitHub (public)",
            "https://github.com/itikelabhaskar/IDBI-Megalodon",
            "Source, tests, model card, connector stubs",
        ),
        (
            "Final product (live prototype)",
            "https://idbi-megalodon.vercel.app",
            "Officer workbench on synthetic data — not a production bank system",
        ),
        (
            "Demo video (3 minutes)",
            "— to be added —",
            "Left blank until the recording URL is ready",
        ),
    ]

    for i, (label, url, note) in enumerate(links):
        top = 2.75 + i * 0.80
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(top), Inches(9.40), Inches(0.70)
        )
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN_SOFT if i < 2 else RGBColor(0xF5, 0xF5, 0xF5)
        shape.line.color.rgb = LINE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        set_run(p0, label, size=11, bold=True, color=GREEN)
        p0.space_after = Pt(1)
        p1 = tf.add_paragraph()
        run = set_run(p1, url, size=12, bold=False, color=RGBColor(0x05, 0x63, 0xC1))
        # Real PowerPoint hyperlink (judges can click from the deck)
        if url.startswith("http"):
            run.hyperlink.address = url
        p1.space_after = Pt(1)
        p2 = tf.add_paragraph()
        set_run(p2, note, size=9, bold=False, color=MUTED)


def patch_slide8_acc_wording(prs: Presentation) -> None:
    """In-place fix if body already exists — avoid full rebuild when only wording changed."""
    slide = prs.slides[7]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "IDBI/ACC" in run.text or "ACC AWS" in run.text:
                    run.text = run.text.replace(
                        "same codebase can target IDBI/ACC AWS path.",
                        "same codebase can target an IDBI-approved AWS path.",
                    ).replace(
                        "IDBI/ACC AWS path",
                        "an IDBI-approved AWS path",
                    )


def patch_slide13_hyperlinks(prs: Presentation) -> None:
    """Ensure GitHub + product URL runs have real hlinkClick targets."""
    slide = prs.slides[12]
    targets = {
        "https://github.com/itikelabhaskar/IDBI-Megalodon",
        "https://idbi-megalodon.vercel.app",
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if text in targets:
                    run.hyperlink.address = text
                    run.font.color.rgb = RGBColor(0x05, 0x63, 0xC1)


def main() -> None:
    prs = Presentation(PATH)
    # Prefer in-place patches for the two Phase-4 review fixes
    patch_slide8_acc_wording(prs)
    # Rebuild slide 13 so hyperlinks are attached at creation time
    fill_slide_13(prs)
    prs.save(PATH)
    print("Phase 4 fixes: Slide 8 ACC wording + Slide 13 clickable hyperlinks.")


if __name__ == "__main__":
    main()
