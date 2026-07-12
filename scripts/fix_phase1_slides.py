"""Phase 1 revision: fix Slide 2/3/4 copy per judge feedback."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PATH = r"D:\IDBI-Megalodon\Prototype Submission Deck _ IDBI Innovate.pptx"


def clear_body_textboxes(slide, min_top_inches: float = 1.15) -> None:
    to_del = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 17:
            top_in = shape.top / 914400
            if top_in > min_top_inches:
                to_del.append(shape)
    for shape in to_del:
        el = shape._element
        el.getparent().remove(el)


def add_body(
    slide,
    text: str,
    left: float = 0.3,
    top: float = 1.5,
    width: float = 9.4,
    height: float = 3.5,
    font_size: int = 13,
    hook_lines: int = 0,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        is_hook = i < hook_lines
        is_section = line in (
            "What makes it different",
            "How it solves Track 03",
            "USP",
            "Proof",
        ) or line.endswith("?") and len(line) < 40
        # section labels: short headings we control
        is_label = line in (
            "What makes it different",
            "How it solves Track 03",
            "USP",
        )
        p.font.size = Pt(15 if is_hook else (12 if is_label else font_size))
        p.font.name = "Calibri"
        p.font.bold = is_hook or is_label
        p.font.color.rgb = (
            RGBColor(0x14, 0x14, 0x14)
            if (is_hook or is_label)
            else RGBColor(0x28, 0x28, 0x28)
        )
        p.space_after = Pt(6 if is_hook else (4 if is_label else 3))
        p.space_before = Pt(6 if is_label and i > 0 else 0)


def shrink_template_prompt(slide) -> None:
    """Turn the official multi-line prompt box into a tiny muted section header strip."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # The opportunities prompt box contains the three template questions
        if "How different" in text and "USP" in text:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Opportunities"
            p.font.size = Pt(16)
            p.font.name = "Calibri"
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            # Shrink the box so body can start higher
            shape.height = Emu(int(0.45 * 914400))
            return


def main() -> None:
    prs = Presentation(PATH)

    for si in (1, 2, 3):
        clear_body_textboxes(prs.slides[si])

    # --- Slide 1 (index 1): Brief — wording fixes ---
    add_body(
        prs.slides[1],
        "Balance sheets are thin. Business behaviour is not.\n"
        "HealthLens turns GST, bank, UPI, EPFO, power and fuel signals into an officer-ready MSME Health Card.\n"
        "\n"
        "• Credit-officer workbench: HealthScore, Go / Conditional / No-Go, limit + IDBI product route, "
        "with reasons an underwriter can defend.\n"
        "• The engine gives reasons and policy logic; the IDBI underwriter owns the decision.\n"
        "• Thin-file NTC/NTB stay in the funnel: when GST or bureau is weak, power and fuel help build "
        "the turnover picture.\n"
        "• Proof now: deployed on synthetic, leakage-aware data; connector stubs are designed to swap "
        "with IDBI / partner sandbox APIs after shortlist.",
        top=1.45,
        height=3.7,
        font_size=12,
        hook_lines=2,
    )

    # --- Slide 2 (index 2): Opportunities — less dense, tiny template label ---
    shrink_template_prompt(prs.slides[2])
    add_body(
        prs.slides[2],
        "What makes it different\n"
        "• Full credit workflow — not a score widget: consent → HealthScore + BRE → maker–checker → CAM.\n"
        "• Power (manufacturing) and fuel (traders/logistics) sit inside the operations sub-score.\n"
        "• Reason codes explain the decision; ML explains only the secondary viability proxy.\n"
        "\n"
        "How it solves Track 03\n"
        "• Credit-invisible MSMEs are not auto-rejected when traditional papers are thin.\n"
        "• GST vs bank vs power/fuel triangulation catches inflated turnover.\n"
        "• Routes to real IDBI products (i-MSME Express, GeM Sahay, CGTMSE) with eligibility gates.\n"
        "\n"
        "USP\n"
        "Pilot-ready decision-support workbench — alternate data → explainable Health Card → "
        "officer Go/No-Go — ULI/OCEN/AA-style stubs ready for sandbox swap.\n"
        "\n"
        "Proof: Frozen MsmeCase contract · leakage guard · audit hash chain · 93 tests",
        top=1.35,
        height=3.9,
        font_size=12,
        hook_lines=0,
    )

    # --- Slide 3 (index 3): Features — unpack last bullet ---
    add_body(
        prs.slides[3],
        "• MSME Financial Health Card — 7 sub-scores → HealthScore + A–D band "
        "(incl. EPFO, power / fuel operations)\n"
        "• Go / Conditional / No-Go with recommended limit — speaks officer language\n"
        "• Business-need detection + IDBI product routing (GeM Sahay gates included)\n"
        "• Fraud & triangulation — GST–bank, power, fuel; Benford advisory only\n"
        "• Explainability + path-to-credit for Refer / Reject — underwriter keeps the call\n"
        "• Maker–checker workflow + printable CAM preview\n"
        "• Live applicant scoring, model card, champion–challenger monitoring and audit trail",
        top=1.45,
        height=3.7,
        font_size=13,
        hook_lines=0,
    )

    prs.save(PATH)
    print("SAVED", PATH)

    prs2 = Presentation(PATH)
    for si in (1, 2, 3):
        print(f"\n=== SLIDE {si + 1} ===")
        for shape in prs2.slides[si].shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print(shape.text_frame.text)
                print("---")
                # char count for slide 3 body
                if si == 2 and "What makes it different" in shape.text_frame.text:
                    print(f"[chars={len(shape.text_frame.text)}]")


if __name__ == "__main__":
    main()
